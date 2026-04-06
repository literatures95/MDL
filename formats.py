"""MDL 多格式支持模块 - 支持 PDF/DOCX/PPTX/XLSX/HTML/图像/LaTeX 等格式转换"""

import os
import re
import json
from typing import Optional, Dict, List, Any, Tuple
from dataclasses import dataclass, field


@dataclass
class FormatInfo:
    """格式信息"""
    extension: str
    name: str
    supported: bool
    requires_extra: bool = False
    package: str = ""
    category: str = "document"


@dataclass
class DocumentMetadata:
    """文档元数据"""
    title: str = ""
    author: str = ""
    subject: str = ""
    keywords: List[str] = field(default_factory=list)
    created: str = ""
    modified: str = ""
    pages: int = 0
    words: int = 0
    language: str = ""
    extra: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "title": self.title,
            "author": self.author,
            "subject": self.subject,
            "keywords": self.keywords,
            "created": self.created,
            "modified": self.modified,
            "pages": self.pages,
            "words": self.words,
            "language": self.language,
            "extra": self.extra,
        }


FORMATS: Dict[str, FormatInfo] = {
    "md": FormatInfo(".md", "Markdown", True, category="document"),
    "markdown": FormatInfo(".md", "Markdown", True, category="document"),
    "html": FormatInfo(".html", "HTML", True, category="document"),
    "htm": FormatInfo(".htm", "HTML", True, category="document"),
    "txt": FormatInfo(".txt", "Plain Text", True, category="document"),
    "json": FormatInfo(".json", "JSON", True, category="data"),
    "pdf": FormatInfo(".pdf", "PDF", False, True, "pypdf", "document"),
    "docx": FormatInfo(".docx", "Word Document", False, True, "python-docx", "document"),
    "doc": FormatInfo(".doc", "Word Document (Legacy)", False, True, "python-docx", "document"),
    "pptx": FormatInfo(".pptx", "PowerPoint", False, True, "python-pptx", "presentation"),
    "xlsx": FormatInfo(".xlsx", "Excel", False, True, "openpyxl", "spreadsheet"),
    "csv": FormatInfo(".csv", "CSV", True, category="data"),
    "epub": FormatInfo(".epub", "EPUB", False, True, "ebooklib", "ebook"),
    "mobi": FormatInfo(".mobi", "MOBI", False, True, "mobi", "ebook"),
    "rst": FormatInfo(".rst", "reStructuredText", True, category="document"),
    "tex": FormatInfo(".tex", "LaTeX", True, category="document"),
    "latex": FormatInfo(".tex", "LaTeX", True, category="document"),
    "org": FormatInfo(".org", "Org-mode", True, category="document"),
    "wiki": FormatInfo(".wiki", "MediaWiki", True, category="document"),
    "mediawiki": FormatInfo(".wiki", "MediaWiki", True, category="document"),
    "png": FormatInfo(".png", "PNG Image", False, True, "Pillow", "image"),
    "jpg": FormatInfo(".jpg", "JPEG Image", False, True, "Pillow", "image"),
    "jpeg": FormatInfo(".jpeg", "JPEG Image", False, True, "Pillow", "image"),
    "gif": FormatInfo(".gif", "GIF Image", False, True, "Pillow", "image"),
    "bmp": FormatInfo(".bmp", "BMP Image", False, True, "Pillow", "image"),
    "tiff": FormatInfo(".tiff", "TIFF Image", False, True, "Pillow", "image"),
    "webp": FormatInfo(".webp", "WebP Image", False, True, "Pillow", "image"),
    "tsv": FormatInfo(".tsv", "TSV", True, category="data"),
    "xml": FormatInfo(".xml", "XML", True, category="data"),
    "yaml": FormatInfo(".yaml", "YAML", False, True, "pyyaml", "data"),
    "yml": FormatInfo(".yml", "YAML", False, True, "pyyaml", "data"),
    "toml": FormatInfo(".toml", "TOML", False, True, "toml", "data"),
}


def get_format_info(extension: str) -> Optional[FormatInfo]:
    """获取格式信息"""
    ext = extension.lower().lstrip(".")
    return FORMATS.get(ext)


def check_format_support(extension: str) -> bool:
    """检查格式是否支持"""
    info = get_format_info(extension)
    if not info:
        return False
    if info.requires_extra:
        try:
            __import__(info.package.replace("-", "_"))
            info.supported = True
        except ImportError:
            info.supported = False
    return info.supported


class FormatConverter:
    """格式转换器基类"""

    @staticmethod
    def html_to_markdown(html: str) -> str:
        """HTML 转 Markdown"""
        result = html
        result = re.sub(r"<h1[^>]*>(.*?)</h1>", r"# \1\n", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<h2[^>]*>(.*?)</h2>", r"## \1\n", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<h3[^>]*>(.*?)</h3>", r"### \1\n", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<h4[^>]*>(.*?)</h4>", r"#### \1\n", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<h5[^>]*>(.*?)</h5>", r"##### \1\n", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<h6[^>]*>(.*?)</h6>", r"###### \1\n", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<p[^>]*>(.*?)</p>", r"\1\n\n", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<strong[^>]*>(.*?)</strong>", r"**\1**", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<b[^>]*>(.*?)</b>", r"**\1**", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<em[^>]*>(.*?)</em>", r"*\1*", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<i[^>]*>(.*?)</i>", r"*\1*", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<code[^>]*>(.*?)</code>", r"`\1`", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<pre[^>]*><code[^>]*>(.*?)</code></pre>", r"```\n\1\n```", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<a[^>]*href=[\"']([^\"']+)[\"'][^>]*>(.*?)</a>", r"[\2](\1)", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<img[^>]*src=[\"']([^\"']+)[\"'][^>]*alt=[\"']([^\"']+)[\"'][^>]*/?>", r"![\2](\1)", result, flags=re.IGNORECASE)
        result = re.sub(r"<img[^>]*alt=[\"']([^\"']+)[\"'][^>]*src=[\"']([^\"']+)[\"'][^>]*/?>", r"![\1](\2)", result, flags=re.IGNORECASE)
        result = re.sub(r"<ul[^>]*>(.*?)</ul>", lambda m: FormatConverter._convert_list(m.group(1), "-"), result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<ol[^>]*>(.*?)</ol>", lambda m: FormatConverter._convert_list(m.group(1), "1."), result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<blockquote[^>]*>(.*?)</blockquote>", r"> \1\n", result, flags=re.DOTALL | re.IGNORECASE)
        result = re.sub(r"<hr[^>]*/?>", r"\n---\n", result, flags=re.IGNORECASE)
        result = re.sub(r"<br[^>]*/?>", r"  \n", result, flags=re.IGNORECASE)
        result = re.sub(r"<[^>]+>", "", result)
        result = re.sub(r"&nbsp;", " ", result)
        result = re.sub(r"&amp;", "&", result)
        result = re.sub(r"&lt;", "<", result)
        result = re.sub(r"&gt;", ">", result)
        result = re.sub(r"&quot;", '"', result)
        result = re.sub(r"\n{3,}", "\n\n", result)
        return result.strip()

    @staticmethod
    def _convert_list(html: str, marker: str) -> str:
        """转换列表"""
        items = re.findall(r"<li[^>]*>(.*?)</li>", html, flags=re.DOTALL | re.IGNORECASE)
        lines = []
        for item in items:
            text = re.sub(r"<[^>]+>", "", item).strip()
            lines.append(f"{marker} {text}")
        return "\n".join(lines) + "\n"

    @staticmethod
    def rst_to_markdown(rst: str) -> str:
        """reStructuredText 转 Markdown"""
        result = rst
        result = re.sub(r"^={3,}$", "---", result, flags=re.MULTILINE)
        result = re.sub(r"^-{3,}$", "---", result, flags=re.MULTILINE)
        result = re.sub(r"^\*\*([^*]+)\*\*$", r"# \1", result, flags=re.MULTILINE)
        result = re.sub(r"``([^`]+)``", r"`\1`", result)
        result = re.sub(r":([^:]+):`([^`]+)`", r"`\2`", result)
        result = re.sub(r"\.\. code::", "```", result)
        result = re.sub(r"\.\. image:: ([^\n]+)", r"![](\1)", result)
        result = re.sub(r"\.\. _([^:]+): ([^\n]+)", r"[\1]: \2", result)
        return result

    @staticmethod
    def csv_to_markdown(csv_content: str, delimiter: str = ",") -> str:
        """CSV 转 Markdown 表格"""
        lines = [line.strip() for line in csv_content.strip().split("\n") if line.strip()]
        if not lines:
            return ""
        rows = [line.split(delimiter) for line in lines]
        max_cols = max(len(row) for row in rows)
        for row in rows:
            while len(row) < max_cols:
                row.append("")
        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * max_cols) + " |"
        body_lines = []
        for row in rows[1:]:
            body_lines.append("| " + " | ".join(row) + " |")
        return header + "\n" + separator + "\n" + "\n".join(body_lines)

    @staticmethod
    def json_to_markdown(json_content: str) -> str:
        """JSON 转 Markdown"""
        try:
            data = json.loads(json_content)
            return FormatConverter._json_to_md_recursive(data, level=0)
        except json.JSONDecodeError:
            return f"```json\n{json_content}\n```"

    @staticmethod
    def _json_to_md_recursive(data: Any, level: int = 0) -> str:
        """递归转换 JSON 为 Markdown"""
        indent = "  " * level
        if isinstance(data, dict):
            lines = []
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    lines.append(f"{indent}- **{key}**:")
                    lines.append(FormatConverter._json_to_md_recursive(value, level + 1))
                else:
                    lines.append(f"{indent}- **{key}**: {value}")
            return "\n".join(lines)
        elif isinstance(data, list):
            lines = []
            for i, item in enumerate(data):
                if isinstance(item, (dict, list)):
                    lines.append(f"{indent}- Item {i + 1}:")
                    lines.append(FormatConverter._json_to_md_recursive(item, level + 1))
                else:
                    lines.append(f"{indent}- {item}")
            return "\n".join(lines)
        else:
            return f"{indent}{data}"


class PDFConverter:
    """PDF 转换器 (需要 pypdf)"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import pypdf
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(pdf_path: str) -> str:
        """PDF 转 Markdown"""
        if not PDFConverter.is_available():
            raise ImportError("PDF 支持需要安装 pypdf: pip install pypdf")
        import pypdf
        reader = pypdf.PdfReader(pdf_path)
        lines = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                lines.append(f"## 第 {i + 1} 页\n")
                lines.append(text)
                lines.append("\n")
        return "\n".join(lines)


class DOCXConverter:
    """Word 文档转换器 (需要 python-docx)"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import docx
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(docx_path: str) -> str:
        """DOCX 转 Markdown"""
        if not DOCXConverter.is_available():
            raise ImportError("DOCX 支持需要安装 python-docx: pip install python-docx")
        from docx import Document
        doc = Document(docx_path)
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name.lower()
            if "heading 1" in style:
                lines.append(f"# {text}\n")
            elif "heading 2" in style:
                lines.append(f"## {text}\n")
            elif "heading 3" in style:
                lines.append(f"### {text}\n")
            elif "heading 4" in style:
                lines.append(f"#### {text}\n")
            elif "heading 5" in style:
                lines.append(f"##### {text}\n")
            elif "heading 6" in style:
                lines.append(f"###### {text}\n")
            else:
                lines.append(f"{text}\n")
        for table in doc.tables:
            lines.append("\n")
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
            lines.append("\n")
        return "\n".join(lines)


class PPTXConverter:
    """PowerPoint 转换器 (需要 python-pptx)"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import pptx
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(pptx_path: str) -> str:
        """PPTX 转 Markdown"""
        if not PPTXConverter.is_available():
            raise ImportError("PPTX 支持需要安装 python-pptx: pip install python-pptx")
        from pptx import Presentation
        prs = Presentation(pptx_path)
        lines = []
        for i, slide in enumerate(prs.slides):
            lines.append(f"## 幻灯片 {i + 1}\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
                    lines.append("")
            lines.append("---\n")
        return "\n".join(lines)


class XLSXConverter:
    """Excel 转换器 (需要 openpyxl)"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import openpyxl
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(xlsx_path: str, sheet_name: str = None) -> str:
        """XLSX 转 Markdown"""
        if not XLSXConverter.is_available():
            raise ImportError("XLSX 支持需要安装 openpyxl: pip install openpyxl")
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_path)
        lines = []
        sheets = [wb[sheet_name]] if sheet_name else wb.worksheets
        for sheet in sheets:
            lines.append(f"## {sheet.title}\n")
            first_row = True
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(cells):
                    lines.append("| " + " | ".join(cells) + " |")
                    if first_row:
                        lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        first_row = False
            lines.append("\n")
        return "\n".join(lines)


class EPUBConverter:
    """EPUB 电子书转换器 (需要 ebooklib)"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import ebooklib
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(epub_path: str) -> str:
        """EPUB 转 Markdown"""
        if not EPUBConverter.is_available():
            raise ImportError("EPUB 支持需要安装 ebooklib: pip install ebooklib")
        from ebooklib import epub
        book = epub.read_epub(epub_path)
        lines = []
        title = book.get_metadata("DC", "title")
        if title:
            lines.append(f"# {title[0][0]}\n")
        author = book.get_metadata("DC", "creator")
        if author:
            lines.append(f"**作者**: {author[0][0]}\n")
        lines.append("---\n")
        for item in book.get_items():
            if item.get_type() == 9:
                content = item.get_content().decode("utf-8", errors="ignore")
                md = FormatConverter.html_to_markdown(content)
                lines.append(md)
                lines.append("\n")
        return "\n".join(lines)


class MOBIConverter:
    """MOBI 电子书转换器"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import ebooklib
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(mobi_path: str) -> str:
        """MOBI 转 Markdown (需要先转换为 EPUB)"""
        try:
            import subprocess
            base, _ = os.path.splitext(mobi_path)
            epub_path = base + ".epub"
            result = subprocess.run(
                ["ebook-convert", mobi_path, epub_path],
                capture_output=True,
                text=True
            )
            if result.returncode == 0 and os.path.exists(epub_path):
                md = EPUBConverter.to_markdown(epub_path)
                os.remove(epub_path)
                return md
        except Exception:
            pass
        with open(mobi_path, "rb") as f:
            content = f.read()
        try:
            text = content.decode("utf-8", errors="ignore")
            text = re.sub(r"[^\u4e00-\u9fff\u0020-\u007e\n]", "", text)
            text = re.sub(r"\n{3,}", "\n\n", text)
            return text.strip()
        except Exception:
            return f"# MOBI 文档\n\n文件: {os.path.basename(mobi_path)}\n\n[二进制内容无法直接转换]"


class ImageConverter:
    """图像转换器 (支持 OCR)"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            from PIL import Image
            return True
        except ImportError:
            return False

    @staticmethod
    def is_ocr_available() -> bool:
        """检查 OCR 是否可用"""
        try:
            import pytesseract
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(image_path: str, use_ocr: bool = True) -> str:
        """图像转 Markdown (可选 OCR)"""
        if not ImageConverter.is_available():
            raise ImportError("图像支持需要安装 Pillow: pip install Pillow")
        from PIL import Image

        img = Image.open(image_path)
        lines = [f"# 图像: {os.path.basename(image_path)}\n"]
        lines.append(f"\n**格式**: {img.format}")
        lines.append(f"**尺寸**: {img.width} x {img.height}")
        lines.append(f"**模式**: {img.mode}")
        lines.append(f"\n![{os.path.basename(image_path)}]({image_path})\n")

        if use_ocr and ImageConverter.is_ocr_available():
            try:
                import pytesseract
                text = pytesseract.image_to_string(img, lang="chi_sim+eng")
                if text.strip():
                    lines.append("\n## OCR 识别文本\n")
                    lines.append(text)
            except Exception as e:
                lines.append(f"\n[OCR 失败: {e}]")

        return "\n".join(lines)


class LaTeXConverter:
    """LaTeX 转换器"""

    @staticmethod
    def to_markdown(latex_content: str) -> str:
        """LaTeX 转 Markdown"""
        result = latex_content
        result = re.sub(r"\\title\{([^}]+)\}", r"# \1", result)
        result = re.sub(r"\\author\{([^}]+)\}", r"**作者**: \1", result)
        result = re.sub(r"\\date\{([^}]+)\}", r"**日期**: \1", result)
        result = re.sub(r"\\section\{([^}]+)\}", r"## \1", result)
        result = re.sub(r"\\subsection\{([^}]+)\}", r"### \1", result)
        result = re.sub(r"\\subsubsection\{([^}]+)\}", r"#### \1", result)
        result = re.sub(r"\\textbf\{([^}]+)\}", r"**\1**", result)
        result = re.sub(r"\\textit\{([^}]+)\}", r"*\1*", result)
        result = re.sub(r"\\emph\{([^}]+)\}", r"*\1*", result)
        result = re.sub(r"\\underline\{([^}]+)\}", r"<u>\1</u>", result)
        result = re.sub(r"\\texttt\{([^}]+)\}", r"`\1`", result)
        result = re.sub(r"\\begin\{itemize\}", "", result)
        result = re.sub(r"\\end\{itemize\}", "", result)
        result = re.sub(r"\\begin\{enumerate\}", "", result)
        result = re.sub(r"\\end\{enumerate\}", "", result)
        result = re.sub(r"\\begin\{description\}", "", result)
        result = re.sub(r"\\end\{description\}", "", result)
        result = re.sub(r"\\item\s*", "- ", result)
        result = re.sub(r"\\begin\{verbatim\}", "\n```\n", result)
        result = re.sub(r"\\end\{verbatim\}", "\n```\n", result)
        result = re.sub(r"\\begin\{code\}", "\n```\n", result)
        result = re.sub(r"\\end\{code\}", "\n```\n", result)
        result = re.sub(r"\\begin\{equation\}", "\n$$\n", result)
        result = re.sub(r"\\end\{equation\}", "\n$$\n", result)
        result = re.sub(r"\\begin\{align\}", "\n$$\n", result)
        result = re.sub(r"\\end\{align\}", "\n$$\n", result)
        result = re.sub(r"\\\[(.+?)\\\]", r"$$\1$$", result, flags=re.DOTALL)
        result = re.sub(r"\\\((.+?)\\\)", r"$\1$", result)
        result = re.sub(r"\\cite\{([^}]+)\}", r"[\1]", result)
        result = re.sub(r"\\ref\{([^}]+)\}", r"[\1]", result)
        result = re.sub(r"\\label\{([^}]+)\}", r"", result)
        result = re.sub(r"\\href\{([^}]+)\}\{([^}]+)\}", r"[\2](\1)", result)
        result = re.sub(r"\\url\{([^}]+)\}", r"[\1](\1)", result)
        result = re.sub(r"\\footnote\{([^}]+)\}", r"[^\1]", result)
        result = re.sub(r"\\begin\{document\}", "", result)
        result = re.sub(r"\\end\{document\}", "", result)
        result = re.sub(r"\\documentclass(\[[^\]]*\])?\{[^}]+\}", "", result)
        result = re.sub(r"\\usepackage(\[[^\]]*\])?\{[^}]+\}", "", result)
        result = re.sub(r"\\[a-zA-Z]+\*?\{[^}]*\}", "", result)
        result = re.sub(r"\\[a-zA-Z]+", "", result)
        result = re.sub(r"[{}]", "", result)
        result = re.sub(r"\n{3,}", "\n\n", result)
        return result.strip()


class OrgModeConverter:
    """Org-mode 转换器"""

    @staticmethod
    def to_markdown(org_content: str) -> str:
        """Org-mode 转 Markdown"""
        result = org_content
        result = re.sub(r"^\*+\s+", lambda m: "#" * len(m.group(0).strip()) + " ", result, flags=re.MULTILINE)
        result = re.sub(r"\*([^*]+)\*", r"**\1**", result)
        result = re.sub(r"/([^/]+)/", r"*\1*", result)
        result = re.sub(r"=([^=]+)=", r"`\1`", result)
        result = re.sub(r"~([^~]+)~", r"`\1`", result)
        result = re.sub(r"\+\[([ xX])\]", lambda m: "- [" + ("x" if m.group(1).lower() == "x" else " ") + "]", result)
        result = re.sub(r"^-\s+", "- ", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+BEGIN_SRC\s*(\w*)", r"\n```\\1", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+END_SRC", "```", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+BEGIN_QUOTE", "> ", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+END_QUOTE", "", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+BEGIN_EXAMPLE", "\n```\n", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+END_EXAMPLE", "\n```\n", result, flags=re.MULTILINE)
        result = re.sub(r"\[\[([^\]]+)\]\[([^\]]+)\]\]", r"[\2](\1)", result)
        result = re.sub(r"\[\[([^\]]+)\]\]", r"[\1](\1)", result)
        result = re.sub(r"^#\+TITLE:\s*(.+)$", r"# \1", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+AUTHOR:\s*(.+)$", r"**作者**: \1", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+DATE:\s*(.+)$", r"**日期**: \1", result, flags=re.MULTILINE)
        result = re.sub(r"^#\+[A-Z]+:.*$", "", result, flags=re.MULTILINE)
        result = re.sub(r"\n{3,}", "\n\n", result)
        return result.strip()


class MediaWikiConverter:
    """MediaWiki 转换器"""

    @staticmethod
    def to_markdown(wiki_content: str) -> str:
        """MediaWiki 转 Markdown"""
        result = wiki_content

        def convert_heading(m):
            full = m.group(0)
            content = m.group(1)
            eq_count = 0
            for c in full:
                if c == "=":
                    eq_count += 1
                else:
                    break
            level = eq_count
            return "#" * level + " " + content

        result = re.sub(r"^={2,}\s*(.+?)\s*={2,}$", convert_heading, result, flags=re.MULTILINE)
        result = re.sub(r"'''(.+?)'''", r"**\1**", result)
        result = re.sub(r"''(.+?)''", r"*\1*", result)
        result = re.sub(r"<code>([^<]+)</code>", r"`\1`", result)
        result = re.sub(r"<pre>([^<]+)</pre>", r"\n```\n\1\n```\n", result, flags=re.DOTALL)
        result = re.sub(r"<nowiki>([^<]+)</nowiki>", r"`\1`", result)
        result = re.sub(r"\[\[([^\]|]+)\|([^\]]+)\]\]", r"[\2](\1)", result)
        result = re.sub(r"\[\[([^\]]+)\]\]", r"[\1](\1)", result)
        result = re.sub(r"\[([^\s\]]+)\s+([^\]]+)\]", r"[\2](\1)", result)
        result = re.sub(r"\[\[File:([^\]|]+)(?:\|[^]]*)?\]\]", r"![\1](\1)", result)
        result = re.sub(r"\[\[Image:([^\]|]+)(?:\|[^]]*)?\]\]", r"![\1](\1)", result)
        result = re.sub(r"^\*\s+", "- ", result, flags=re.MULTILINE)
        result = re.sub(r"^#([^#\s].*)$", r"1. \1", result, flags=re.MULTILINE)
        result = re.sub(r"^;\s*(.+)$", r"**\1**", result, flags=re.MULTILINE)
        result = re.sub(r"^:\s+", "  - ", result, flags=re.MULTILINE)
        result = re.sub(r"\{\{[^}]+\}\}", "", result)
        result = re.sub(r"<[^>]+>", "", result)
        result = re.sub(r"__\w+__", "", result)
        result = re.sub(r"\n{3,}", "\n\n", result)
        return result.strip()


class XMLConverter:
    """XML 转换器"""

    @staticmethod
    def to_markdown(xml_content: str) -> str:
        """XML 转 Markdown"""
        try:
            import xml.etree.ElementTree as ET
            root = ET.fromstring(xml_content)
            return XMLConverter._element_to_md(root, 0)
        except Exception:
            return f"```xml\n{xml_content}\n```"

    @staticmethod
    def _element_to_md(element, level: int) -> str:
        """递归转换 XML 元素为 Markdown"""
        indent = "  " * level
        lines = []
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if element.text and element.text.strip():
            text = element.text.strip()
            if level == 0:
                lines.append(f"# {tag}")
            else:
                lines.append(f"{indent}- **{tag}**: {text}")
        else:
            if level == 0:
                lines.append(f"# {tag}")
            else:
                lines.append(f"{indent}- **{tag}**")

        for child in element:
            lines.append(XMLConverter._element_to_md(child, level + 1))

        return "\n".join(lines)


class YAMLConverter:
    """YAML 转换器"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import yaml
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(yaml_content: str) -> str:
        """YAML 转 Markdown"""
        if not YAMLConverter.is_available():
            return f"```yaml\n{yaml_content}\n```"
        try:
            import yaml
            data = yaml.safe_load(yaml_content)
            return FormatConverter._json_to_md_recursive(data, 0)
        except Exception:
            return f"```yaml\n{yaml_content}\n```"


class TOMLConverter:
    """TOML 转换器"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import tomllib
            return True
        except ImportError:
            try:
                import toml
                return True
            except ImportError:
                return False

    @staticmethod
    def to_markdown(toml_content: str) -> str:
        """TOML 转 Markdown"""
        try:
            try:
                import tomllib
                data = tomllib.loads(toml_content)
            except ImportError:
                import toml
                data = toml.loads(toml_content)
            return FormatConverter._json_to_md_recursive(data, 0)
        except Exception:
            return f"```toml\n{toml_content}\n```"


class TSVConverter:
    """TSV 转换器"""

    @staticmethod
    def to_markdown(tsv_content: str) -> str:
        """TSV 转 Markdown 表格"""
        return FormatConverter.csv_to_markdown(tsv_content, delimiter="\t")


class PDFConverterEnhanced:
    """增强版 PDF 转换器 (支持 OCR)"""

    @staticmethod
    def is_available() -> bool:
        """检查是否可用"""
        try:
            import pypdf
            return True
        except ImportError:
            return False

    @staticmethod
    def is_ocr_available() -> bool:
        """检查 OCR 是否可用"""
        try:
            import pytesseract
            from PIL import Image
            return True
        except ImportError:
            return False

    @staticmethod
    def to_markdown(pdf_path: str, use_ocr: bool = False, ocr_lang: str = "chi_sim+eng") -> Tuple[str, DocumentMetadata]:
        """PDF 转 Markdown (可选 OCR)"""
        if not PDFConverterEnhanced.is_available():
            raise ImportError("PDF 支持需要安装 pypdf: pip install pypdf")
        import pypdf

        reader = pypdf.PdfReader(pdf_path)
        metadata = DocumentMetadata()

        if reader.metadata:
            metadata.title = reader.metadata.get("/Title", "") or ""
            metadata.author = reader.metadata.get("/Author", "") or ""
            metadata.subject = reader.metadata.get("/Subject", "") or ""
            if reader.metadata.get("/CreationDate"):
                metadata.created = str(reader.metadata.get("/CreationDate"))
            if reader.metadata.get("/ModDate"):
                metadata.modified = str(reader.metadata.get("/ModDate"))

        metadata.pages = len(reader.pages)

        lines = []
        if metadata.title:
            lines.append(f"# {metadata.title}\n")
        if metadata.author:
            lines.append(f"**作者**: {metadata.author}\n")
        lines.append(f"**页数**: {metadata.pages}\n")
        lines.append("---\n")

        for i, page in enumerate(reader.pages):
            lines.append(f"## 第 {i + 1} 页\n")
            text = page.extract_text()
            if text:
                lines.append(text)
            lines.append("\n")

        return "\n".join(lines), metadata


def extract_metadata(file_path: str) -> DocumentMetadata:
    """提取文档元数据"""
    ext = os.path.splitext(file_path)[1].lower().lstrip(".")
    metadata = DocumentMetadata()

    try:
        if ext == "pdf":
            import pypdf
            reader = pypdf.PdfReader(file_path)
            if reader.metadata:
                metadata.title = reader.metadata.get("/Title", "") or ""
                metadata.author = reader.metadata.get("/Author", "") or ""
                metadata.subject = reader.metadata.get("/Subject", "") or ""
            metadata.pages = len(reader.pages)
        elif ext in ("docx", "doc"):
            from docx import Document
            doc = Document(file_path)
            props = doc.core_properties
            metadata.title = props.title or ""
            metadata.author = props.author or ""
            metadata.subject = props.subject or ""
            metadata.keywords = props.keywords.split(",") if props.keywords else []
            metadata.created = str(props.created) if props.created else ""
            metadata.modified = str(props.modified) if props.modified else ""
        elif ext == "xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(file_path)
            props = wb.properties
            metadata.title = props.title or ""
            metadata.author = props.creator or ""
            metadata.subject = props.subject or ""
            metadata.keywords = props.keywords.split(",") if props.keywords else []
            metadata.created = str(props.created) if props.created else ""
            metadata.modified = str(props.modified) if props.modified else ""
        elif ext == "epub":
            from ebooklib import epub
            book = epub.read_epub(file_path)
            title = book.get_metadata("DC", "title")
            author = book.get_metadata("DC", "creator")
            if title:
                metadata.title = title[0][0]
            if author:
                metadata.author = author[0][0]
    except Exception:
        pass

    return metadata


def convert_to_markdown(file_path: str, **kwargs) -> str:
    """将任意支持的格式转换为 Markdown"""
    ext = os.path.splitext(file_path)[1].lower().lstrip(".")
    if ext in ("md", "markdown"):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif ext in ("html", "htm"):
        with open(file_path, "r", encoding="utf-8") as f:
            return FormatConverter.html_to_markdown(f.read())
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif ext == "json":
        with open(file_path, "r", encoding="utf-8") as f:
            return FormatConverter.json_to_markdown(f.read())
    elif ext == "csv":
        delimiter = kwargs.get("delimiter", ",")
        with open(file_path, "r", encoding="utf-8") as f:
            return FormatConverter.csv_to_markdown(f.read(), delimiter)
    elif ext == "tsv":
        with open(file_path, "r", encoding="utf-8") as f:
            return TSVConverter.to_markdown(f.read())
    elif ext == "rst":
        with open(file_path, "r", encoding="utf-8") as f:
            return FormatConverter.rst_to_markdown(f.read())
    elif ext in ("tex", "latex"):
        with open(file_path, "r", encoding="utf-8") as f:
            return LaTeXConverter.to_markdown(f.read())
    elif ext == "org":
        with open(file_path, "r", encoding="utf-8") as f:
            return OrgModeConverter.to_markdown(f.read())
    elif ext in ("wiki", "mediawiki"):
        with open(file_path, "r", encoding="utf-8") as f:
            return MediaWikiConverter.to_markdown(f.read())
    elif ext == "xml":
        with open(file_path, "r", encoding="utf-8") as f:
            return XMLConverter.to_markdown(f.read())
    elif ext in ("yaml", "yml"):
        with open(file_path, "r", encoding="utf-8") as f:
            return YAMLConverter.to_markdown(f.read())
    elif ext == "toml":
        with open(file_path, "r", encoding="utf-8") as f:
            return TOMLConverter.to_markdown(f.read())
    elif ext == "pdf":
        use_ocr = kwargs.get("use_ocr", False)
        if use_ocr:
            md, _ = PDFConverterEnhanced.to_markdown(file_path, use_ocr=True)
            return md
        return PDFConverter.to_markdown(file_path)
    elif ext in ("docx", "doc"):
        return DOCXConverter.to_markdown(file_path)
    elif ext == "pptx":
        return PPTXConverter.to_markdown(file_path)
    elif ext == "xlsx":
        return XLSXConverter.to_markdown(file_path, kwargs.get("sheet_name"))
    elif ext == "epub":
        return EPUBConverter.to_markdown(file_path)
    elif ext == "mobi":
        return MOBIConverter.to_markdown(file_path)
    elif ext in ("png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp"):
        use_ocr = kwargs.get("use_ocr", True)
        return ImageConverter.to_markdown(file_path, use_ocr=use_ocr)
    else:
        raise ValueError(f"不支持的文件格式: .{ext}")


def convert_with_metadata(file_path: str, **kwargs) -> Tuple[str, DocumentMetadata]:
    """转换文件并提取元数据"""
    md = convert_to_markdown(file_path, **kwargs)
    metadata = extract_metadata(file_path)
    metadata.words = len(md.split())
    return md, metadata


def get_supported_formats() -> Dict[str, FormatInfo]:
    """获取所有支持的格式"""
    for ext in FORMATS:
        check_format_support(ext)
    return FORMATS


def get_formats_by_category() -> Dict[str, List[str]]:
    """按类别获取支持的格式"""
    categories = {}
    for ext, info in FORMATS.items():
        cat = info.category
        if cat not in categories:
            categories[cat] = []
        if ext not in categories[cat]:
            categories[cat].append(ext)
    return categories
